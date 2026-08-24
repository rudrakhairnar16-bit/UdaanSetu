"""Fill the official SIH 2026 template with complete UdaanSetu content, preserving layout."""
from pptx import Presentation
from pptx.util import Pt

SRC = r"C:\Users\Rudra\Desktop\UdaanSetu\SIH_Submission\SIH2026-Official-Template.pptx"
OUT = r"C:\Users\Rudra\Desktop\UdaanSetu\SIH_Submission\UdaanSetu_SIH2026_Demo.pptx"

prs = Presentation(SRC)
slides = list(prs.slides)


def find(slide, sid):
    for sh in slide.shapes:
        if sh.shape_id == sid:
            return sh
    return None


def fill_tf(tf, lines, default_size=28, default_bold=None):
    for p in list(tf.paragraphs):
        p._p.getparent().remove(p._p)
    tf.word_wrap = True
    for text, level, bold, size in lines:
        p = tf.add_paragraph()
        p.level = level
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = "Arial"
        f.size = Pt(size if size else default_size)
        f.bold = bold if bold is not None else default_bold


def fill_textbox(slide, sid, lines, default_size=28, default_bold=None):
    sh = find(slide, sid)
    if sh is not None and sh.has_text_frame:
        fill_tf(sh.text_frame, lines, default_size, default_bold)


def grow_textbox(slide, sid, bottom_target=5900000):
    """Extend a textbox so it uses the space before the footer."""
    sh = find(slide, sid)
    if sh is not None and sh.has_text_frame:
        sh.height = bottom_target - sh.top
        sh.text_frame.auto_size = 1  # MSO_AUTO_SIZE.SHRINK_TEXT_ON_OVERFLOW


# ============ Slide 1 - TITLE PAGE ============
s1 = slides[0]
fill_textbox(s1, 10, [
    ("Problem Statement ID \u2013 SIH1608", 0, False, 26),
    ("Problem Statement Title \u2013 Create an Innovation Ecosystem Platform", 0, False, 26),
    ("Theme \u2013 Miscellaneous", 0, False, 26),
    ("PS Category \u2013 Software", 0, False, 26),
    ("Team ID \u2013 [As Per Portal]", 0, False, 26),
    ("Team Name \u2013 [TEAM NAME] (Registered on Portal)", 0, False, 26),
], default_size=26)

# ============ Slide 2 - PROPOSED SOLUTION ============
s2 = slides[1]
fill_textbox(s2, 15361, [
    ("UdaanSetu", 0, True, 40),
    ("Bridge to Flight \u2014 Research to Impact", 0, False, 24),
], default_size=40, default_bold=False)
fill_textbox(s2, 15362, [
    ("Proposed Solution (Describe your Idea/Solution/Prototype)", 0, True, 22),
    ("Detailed explanation of the proposed solution", 0, True, 16),
    ("\u2022 Unified innovation-ecosystem platform guiding each project end-to-end: research \u2192 patent \u2192 mentorship \u2192 funding \u2192 incubation \u2192 startup \u2192 impact.", 0, False, 16),
    ("\u2022 Role-based dashboards for Researchers, Mentors, Investors, Incubators & Admins.", 0, False, 16),
    ("\u2022 ML-driven risk prediction, semantic duplicate detection and smart matching.", 0, False, 16),
    ("\u2022 Government integrations: Aadhaar eKYC, DigiLocker, Startup India, IP India, ONDC.", 0, False, 16),
    ("How it addresses the problem", 0, True, 16),
    ("\u2022 Removes fragmentation: every stakeholder operates on one living pipeline.", 0, False, 16),
    ("\u2022 Makes schemes, mentors and funding visible to the right innovator automatically.", 0, False, 16),
    ("\u2022 Transparent IPR tracking with audit trails for all ecosystem actors.", 0, False, 16),
    ("Innovation and uniqueness of the solution", 0, True, 16),
    ("\u2022 Complete lab-to-market lifecycle on a single platform (first of its kind).", 0, False, 16),
    ("\u2022 Explainable AI risk scoring + retraining on real data.", 0, False, 16),
    ("\u2022 Production-grade: 60+ REST APIs, Dockerized, 190+ automated tests, security-first.", 0, False, 16),
], default_size=16)
tn = find(s2, 10)
if tn is not None:
    fill_tf(tn.text_frame, [("[TEAM NAME]", 0, True, 20)], default_size=20)
grow_textbox(s2, 15362, 5900000)

# ============ Slide 3 - TECHNICAL APPROACH ============
s3 = slides[2]
fill_textbox(s3, 17410, [
    ("Technologies to be used", 0, True, 20),
    ("\u2022 Frontend: Next.js 15 \u00b7 React 19 \u00b7 TypeScript \u00b7 accessible design system (WCAG 2.1 AA)", 0, False, 15),
    ("\u2022 Backend: FastAPI \u00b7 Python \u00b7 SQLAlchemy 2 \u00b7 Pydantic v2 \u00b7 Uvicorn", 0, False, 15),
    ("\u2022 Database: PostgreSQL + Redis cache (in-memory fallback)", 0, False, 15),
    ("\u2022 AI/ML: scikit-learn (GBM, clustering) \u00b7 sentence-transformers embeddings", 0, False, 15),
    ("\u2022 Security: JWT \u00b7 Argon2 \u00b7 RBAC \u00b7 rate limiting \u00b7 audit log", 0, False, 15),
    ("\u2022 Infra: Docker Compose \u00b7 health checks \u00b7 CI-ready test suites", 0, False, 15),
    ("Methodology and process for implementation", 0, True, 20),
    ("\u2022 Lifecycle model: Research \u2192 Innovation \u2192 IPR \u2192 Mentor/Funding/Incubator \u2192 Startup \u2192 Impact", 0, False, 15),
    ("\u2022 Agile sprints; iterative demos; milestone-driven delivery.", 0, False, 15),
    ("\u2022 Working prototype: 60+ REST endpoints, live dashboards, Docker Compose full-stack demo.", 0, False, 15),
    ("\u2022 ML pipeline: feature extraction \u2192 training \u2192 registry \u2192 serving \u2192 retrain on real data.", 0, False, 15),
], default_size=15)
tn = find(s3, 11)
if tn is not None:
    fill_tf(tn.text_frame, [("[TEAM NAME]", 0, True, 20)], default_size=20)
grow_textbox(s3, 17410, 5900000)

# ============ Slide 4 - FEASIBILITY AND VIABILITY ============
s4 = slides[3]
fill_textbox(s4, 17410, [
    ("Analysis of the feasibility of the idea", 0, True, 18),
    ("\u2022 Leverages existing govt infrastructure: Aadhaar, DigiLocker, Startup India, IP India, ONDC APIs.", 0, False, 14),
    ("\u2022 Technically viable: proven, mature stack; scalable cloud deployment.", 0, False, 14),
    ("\u2022 Strong demand: fragmented innovation ecosystem across Indian institutes & districts.", 0, False, 14),
    ("\u2022 Operationally viable: demo-ready with realistic seeded data (127 projects).", 0, False, 14),
    ("Potential challenges and risks", 0, True, 18),
    ("\u2022 Data privacy & consent for government data.", 0, False, 14),
    ("\u2022 Adoption across institutes and government departments.", 0, False, 14),
    ("\u2022 Model quality on sparse early-stage data.", 0, False, 14),
    ("Strategies for overcoming these challenges", 0, True, 18),
    ("\u2022 RBAC, audit logs and consent-first data design.", 0, False, 14),
    ("\u2022 Institute onboarding, training & multilingual-friendly UI.", 0, False, 14),
    ("\u2022 Synthetic-data fallback + retraining on real records when samples grow.", 0, False, 14),
], default_size=14)
tn = find(s4, 12)
if tn is not None:
    fill_tf(tn.text_frame, [("[TEAM NAME]", 0, True, 20)], default_size=20)
grow_textbox(s4, 17410, 5900000)

# ============ Slide 5 - IMPACT AND BENEFITS ============
s5 = slides[4]
fill_textbox(s5, 17410, [
    ("Potential impact on the target audience", 0, True, 18),
    ("\u2022 Researchers: guided, measurable path from paper to product.", 0, False, 14),
    ("\u2022 Innovators/Startups: automatic access to mentors, schemes and funding.", 0, False, 14),
    ("\u2022 Investors/Incubators: visible, risk-scored pipeline to pick winners.", 0, False, 14),
    ("\u2022 Government: real-time transparency across the national innovation pipeline.", 0, False, 14),
    ("Benefits of the solution (social, economic, environmental)", 0, True, 18),
    ("\u2022 Social: democratizes innovation access across districts and demographics.", 0, False, 14),
    ("\u2022 Economic: shortens lab-to-market time (3\u20135 yrs \u2192 months), creates startups & jobs.", 0, False, 14),
    ("\u2022 Environmental: paperless IPR and scheme workflows reduce administrative waste.", 0, False, 14),
], default_size=14)
tn = find(s5, 12)
if tn is not None:
    fill_tf(tn.text_frame, [("[TEAM NAME]", 0, True, 20)], default_size=20)
grow_textbox(s5, 17410, 5900000)

# ============ Slide 6 - RESEARCH AND REFERENCES ============
s6 = slides[5]
fill_textbox(s6, 17410, [
    ("\u2022 India publishes 50,000+ research papers/year; fewer than 500 become startups.", 0, False, 16),
    ("\u2022 NITI Aayog \u2013 Innovation & Entrepreneurship ecosystem reports.", 0, False, 16),
    ("\u2022 Startup India / DPIIT startup policy and scheme documentation.", 0, False, 16),
    ("\u2022 AICTE / MoE Innovation Cell \u2013 Smart India Hackathon 2026 guidelines.", 0, False, 16),
    ("\u2022 Technology-transfer and innovation-ecosystem academic literature.", 0, False, 16),
    ("\u2022 Live demo & repository: github.com/rudrakhairnar16-bit/UdaanSetu", 0, False, 16),
], default_size=16)
tn = find(s6, 9)
if tn is not None:
    fill_tf(tn.text_frame, [("[TEAM NAME]", 0, True, 20)], default_size=20)
grow_textbox(s6, 17410, 5900000)

prs.save(OUT)
print("Saved:", OUT)