"""Generate UdaanSetu SIH 2026 demo PPT."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand palette
DARK = RGBColor(0x0C, 0x3B, 0x26)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
MID = RGBColor(0x06, 0x4E, 0x3B)
SOFT = RGBColor(0xF0, 0xFD, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x4B, 0x55, 0x63)
LIGHT = RGBColor(0xF9, 0xFA, 0xFB)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, fill, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def txt(s, x, y, w, h, runs, size=14, color=GRAY, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [(runs, {})]
    for i, (text, opts) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = opts.get('spacing', spacing)
        r = p.add_run()
        r.text = text
        f = r.font
        f.size = Pt(opts.get('size', size))
        f.bold = opts.get('bold', bold)
        f.color.rgb = opts.get('color', color)
        f.name = 'Calibri'
    return tb


def header(s, title, sub=None):
    rect(s, 0, 0, SW, Inches(1.0), DARK)
    rect(s, 0, Inches(1.0), SW, Pt(3), GREEN)
    txt(s, Inches(0.6), Inches(0.18), Inches(9), Inches(0.6),
        [(title, {'size': 26, 'bold': True, 'color': WHITE})])
    if sub:
        txt(s, Inches(9.2), Inches(0.3), Inches(3.6), Inches(0.5),
            [(sub, {'size': 12, 'color': RGBColor(0xA7, 0xF3, 0xD0), 'align': PP_ALIGN.RIGHT})])


def bullets(s, x, y, w, h, items, size=16, gap=0.9, color=GRAY):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = gap
        p.space_after = Pt(4)
        if isinstance(it, tuple):
            head, body = it
            r = p.add_run(); r.text = "•  "
            r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = GREEN
            r1 = p.add_run(); r1.text = head
            r1.font.size = Pt(size); r1.font.bold = True; r1.font.color.rgb = DARK
            r2 = p.add_run(); r2.text = " — " + body
            r2.font.size = Pt(size); r2.font.color.rgb = color
        else:
            r = p.add_run(); r.text = "•  " + it
            r.font.size = Pt(size); r.font.color.rgb = color
    return tb


def footer(s, n):
    rect(s, 0, SH - Inches(0.45), SW, Inches(0.45), LIGHT)
    txt(s, Inches(0.6), SH - Inches(0.36), Inches(8), Inches(0.3),
        [("UdaanSetu — SIH 2026 | Problem ID SIH1608",
          {'size': 10, 'color': GRAY})])
    txt(s, Inches(12.2), SH - Inches(0.36), Inches(0.6), Inches(0.3),
        [(str(n), {'size': 10, 'color': GRAY, 'align': PP_ALIGN.RIGHT})])


# ---------------------------------------------------------------- Slide 1 - Title
s = slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, Inches(4.7), SW, Pt(3), GREEN)
rect(s, Inches(10.9), 0, Inches(2.4), SH, RGBColor(0x0A, 0x2F, 0x1E))
txt(s, Inches(0.9), Inches(1.1), Inches(9), Inches(1.0),
    [("↗ UdaanSetu", {'size': 60, 'bold': True, 'color': WHITE})])
txt(s, Inches(0.95), Inches(2.0), Inches(9), Inches(0.6),
    [("Bridge to Flight — Research to Impact", {'size': 24, 'color': RGBColor(0x86, 0xEF, 0xAC)})])
txt(s, Inches(0.95), Inches(2.7), Inches(9.5), Inches(1.6),
    [("An Innovation Ecosystem Platform that takes research from the lab "
      "to market-ready startups, guided at every step.", {'size': 18, 'color': RGBColor(0xD1, 0xFA, 0xE5)})])
txt(s, Inches(0.95), Inches(5.1), Inches(9), Inches(1.6), [
    ("Smart India Hackathon 2026", {'size': 20, 'bold': True, 'color': WHITE}),
    ("Problem ID: SIH1608  •  Track: Innovation Ecosystem", {'size': 15, 'color': RGBColor(0xA7, 0xF3, 0xD0)}),
    ("Team Name: [TEAM NAME]  •  Institute: [INSTITUTE]", {'size': 14, 'color': RGBColor(0x86, 0xEF, 0xAC)}),
])

# ---------------------------------------------------------------- Slide 2 - Problem
s = slide()
header(s, "The Problem", "SIH1608")
txt(s, Inches(0.6), Inches(1.35), Inches(12), Inches(0.5),
    [("India produces world-class research — but it rarely becomes products.",
      {'size': 20, 'bold': True, 'color': DARK})])
bullets(s, Inches(0.6), Inches(2.0), Inches(6.2), Inches(4.6), [
    ("Fragmented journey", "research → patent → funding → startup has no single guide"),
    ("Invisible schemes", "government grants & schemes never reach innovators"),
    ("Opaque IPR", "patent filing is confusing, slow, and untracked"),
    ("No mentorship", "founders don't know who to ask or what step comes next"),
    ("No ecosystem view", "stakeholders can't see the full innovation pipeline"),
], size=16)
rect(s, Inches(7.2), Inches(2.0), Inches(5.5), Inches(4.4), SOFT, line=RGBColor(0xBB, 0xF7, 0xD0))
txt(s, Inches(7.6), Inches(2.2), Inches(4.8), Inches(0.4),
    [("The Cost", {'size': 18, 'bold': True, 'color': MID})])
bullets(s, Inches(7.6), Inches(2.7), Inches(4.8), Inches(3.6), [
    "50,000+ research papers published in India every year",
    "< 500 ever become startups",
    "95% of research never reaches the market",
    "Average lab-to-market time: 3–5 years",
], size=14, gap=1.1)
footer(s, 2)

# ---------------------------------------------------------------- Slide 3 - Solution
s = slide()
header(s, "Our Solution", "SIH1608")
rect(s, Inches(0.6), Inches(1.35), Inches(12.1), Inches(1.15), SOFT, line=RGBColor(0xBB, 0xF7, 0xD0))
txt(s, Inches(1.0), Inches(1.5), Inches(11.2), Inches(0.8),
    [("UdaanSetu unifies the entire innovation lifecycle on one platform — "
      "guiding every project from research to real-world impact.",
      {'size': 18, 'bold': True, 'color': MID})])
steps = [
    ("1", "Research", "Capture projects, progress & milestones"),
    ("2", "Innovation", "Semantic duplicate detection & AI recommendations"),
    ("3", "IPR / Patent", "Track filings, stages & reminders"),
    ("4", "Mentor · Funding · Incubator", "AI-powered smart matching"),
    ("5", "Startup", "Launch, track jobs & revenue"),
    ("6", "Impact", "Measure sector & district outcomes"),
]
x = Inches(0.6)
for num, title, desc in steps:
    card = rect(s, x, Inches(2.75), Inches(1.9), Inches(2.9), WHITE, line=RGBColor(0xE5, 0xE7, 0xEB))
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.65), Inches(3.0), Inches(0.6), Inches(0.6))
    circ.fill.solid(); circ.fill.fore_color.rgb = GREEN; circ.line.fill.background()
    circ.shadow.inherit = False
    tf = circ.text_frame; tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE
    txt(s, x + Inches(0.1), Inches(3.75), Inches(1.7), Inches(0.9),
        [(title, {'size': 14, 'bold': True, 'color': DARK, 'align': PP_ALIGN.CENTER})], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + Inches(0.15), Inches(4.6), Inches(1.6), Inches(1.0),
        [(desc, {'size': 11, 'color': GRAY, 'align': PP_ALIGN.CENTER})])
    x += Inches(2.04)
footer(s, 3)

# ---------------------------------------------------------------- Slide 4 - Features
s = slide()
header(s, "Key Features", "Live Demo-Ready")
feats = [
    ("Role-Based Dashboards", "Admin · Researcher · Mentor · Investor · Incubator — each with tailored views"),
    ("ML Risk Prediction", "Gradient-boosting model flags at-risk projects with explainable reasons"),
    ("Semantic Search & Matching", "sentence-transformer embeddings connect innovators, mentors & schemes"),
    ("Duplicate Detection", "NLP clustering prevents redundant projects in the ecosystem"),
    ("Government Integrations", "Aadhaar, DigiLocker, Startup India, IP India, ONDC (mock, demo mode)"),
    ("Security First", "JWT auth, Argon2 hashing, RBAC, rate-limiting, audit log"),
]
x, y = Inches(0.6), Inches(1.4)
for title, desc in feats:
    card = rect(s, x, y, Inches(5.9), Inches(1.5), WHITE, line=RGBColor(0xE5, 0xE7, 0xEB))
    rect(s, x, y, Pt(4), Inches(1.5), GREEN)
    txt(s, x + Inches(0.3), y + Inches(0.15), Inches(5.4), Inches(0.4),
        [(title, {'size': 15, 'bold': True, 'color': DARK})])
    txt(s, x + Inches(0.3), y + Inches(0.6), Inches(5.4), Inches(0.8),
        [(desc, {'size': 12, 'color': GRAY})])
    if x < Inches(6.8):
        x += Inches(6.2)
    else:
        x, y = Inches(0.6), y + Inches(1.7)
footer(s, 4)

# ---------------------------------------------------------------- Slide 5 - Tech Stack
s = slide()
header(s, "Technology Stack", "Production-Grade")
tech = [
    ("Frontend", "Next.js 15 · React 19 · TypeScript · CSS Design System (WCAG AA)"),
    ("Backend", "FastAPI · Python 3.12 · SQLAlchemy · Pydantic v2"),
    ("Database", "PostgreSQL · Redis caching layer"),
    ("ML / AI", "scikit-learn · sentence-transformers · real-data training pipeline"),
    ("Security", "JWT · Argon2 · RBAC · input sanitization · rate limiting"),
    ("Infra", "Docker Compose · automated tests · health checks · model registry"),
]
y = Inches(1.4)
for name, desc in tech:
    rect(s, Inches(0.6), y, Inches(3.2), Inches(0.8), DARK)
    txt(s, Inches(0.7), y + Inches(0.2), Inches(3.0), Inches(0.4),
        [(name, {'size': 15, 'bold': True, 'color': WHITE})])
    rect(s, Inches(3.9), y, Inches(8.8), Inches(0.8), LIGHT, line=RGBColor(0xE5, 0xE7, 0xEB))
    txt(s, Inches(4.1), y + Inches(0.18), Inches(8.5), Inches(0.5),
        [(desc, {'size': 13, 'color': GRAY})], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.98)
txt(s, Inches(0.6), Inches(6.9), Inches(12), Inches(0.4),
    [("172 backend tests + 19 frontend tests passing · Dockerized full-stack demo",
      {'size': 12, 'bold': True, 'color': GREEN})])
footer(s, 5)

# ---------------------------------------------------------------- Slide 6 - Architecture
s = slide()
header(s, "System Architecture", "Clean, Scalable")
# Draw blocks
def block(x, y, w, h, title, sub, fill, tcolor=WHITE):
    b = rect(s, x, y, w, h, fill)
    txt(s, x + Inches(0.15), y + Inches(0.12), w - Inches(0.3), Inches(0.35),
        [(title, {'size': 13, 'bold': True, 'color': tcolor})])
    txt(s, x + Inches(0.15), y + Inches(0.5), w - Inches(0.3), Inches(0.5),
        [(sub, {'size': 10.5, 'color': RGBColor(0xD1, 0xFA, 0xE5) if tcolor == WHITE else GRAY})])
    return b

block(Inches(0.7), Inches(1.5), Inches(3.4), Inches(1.2), "React Frontend", "Next.js · TypeScript", DARK)
block(Inches(0.7), Inches(3.0), Inches(3.4), Inches(1.2), "FastAPI Backend", "60+ REST endpoints", MID)
block(Inches(0.7), Inches(4.5), Inches(3.4), Inches(1.2), "PostgreSQL", "Normalized schema", RGBColor(0x0A, 0x2F, 0x1E))
block(Inches(5.2), Inches(1.5), Inches(3.4), Inches(1.2), "ML Engine", "Risk · Semantic · Matching", GREEN)
block(Inches(5.2), Inches(3.0), Inches(3.4), Inches(1.2), "Govt. Integrations", "Aadhaar · DigiLocker · ONDC", GREEN)
block(Inches(5.2), Inches(4.5), Inches(3.4), Inches(1.2), "Redis Cache", "Dashboard & analytics", MID)
block(Inches(9.7), Inches(1.5), Inches(3.0), Inches(1.2), "Audit Log", "Every action tracked", AMBER, DARK)
block(Inches(9.7), Inches(3.0), Inches(3.0), Inches(1.2), "Notifications", "Real-time updates", AMBER, DARK)
block(Inches(9.7), Inches(4.5), Inches(3.0), Inches(1.2), "Auth & RBAC", "JWT · Argon2", AMBER, DARK)
txt(s, Inches(0.7), Inches(6.0), Inches(12), Inches(1.0),
    [("Stateless APIs · Containerized with Docker Compose · Health-checked services · "
      "Caching via Redis with in-memory fallback", {'size': 12, 'color': GRAY})])
footer(s, 6)

# ---------------------------------------------------------------- Slide 7 - ML/AI
s = slide()
header(s, "AI / ML Capabilities", "Explainable & Real-Data Ready")
rows = [
    ("Risk Prediction", "Gradient-boosting model scores every project (0–100) with clear reasons: overdue milestones, low progress, stalled stage."),
    ("Smart Matching", "sentence-transformers embeddings match innovators ↔ mentors, schemes, incubators by semantic similarity."),
    ("Duplicate Detection", "Agglomerative clustering on embeddings flags near-duplicate research / IP to keep the ecosystem clean."),
    ("Success Prediction", "Probability + confidence interval for a project's commercialization, with comparable project insights."),
    ("Retraining Pipeline", "Trains on real DB records when ≥20 samples exist; synthetic fallback otherwise. Model registry + versioning."),
]
y = Inches(1.4)
for title, desc in rows:
    rect(s, Inches(0.6), y, Inches(12.1), Inches(1.0), WHITE, line=RGBColor(0xE5, 0xE7, 0xEB))
    rect(s, Inches(0.6), y, Inches(2.7), Inches(1.0), SOFT, line=RGBColor(0xBB, 0xF7, 0xD0))
    txt(s, Inches(0.75), y + Inches(0.25), Inches(2.5), Inches(0.5),
        [(title, {'size': 14, 'bold': True, 'color': MID})])
    txt(s, Inches(3.5), y + Inches(0.12), Inches(9.0), Inches(0.8),
        [(desc, {'size': 12.5, 'color': GRAY})], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.12)
footer(s, 7)

# ---------------------------------------------------------------- Slide 8 - Demo Flow
s = slide()
header(s, "Live Demo Walkthrough", "5-Minute Script")
steps = [
    ("00:00", "Login as Researcher", "admin@udaansetu.demo / Demo@123"),
    ("00:30", "Create Research Project", "set stage, district, sector, funding need"),
    ("01:15", "Add Milestones", "track progress; overdue flagged red"),
    ("01:45", "AI Risk Assessment", "score + reasons; at-risk project highlighted"),
    ("02:15", "Smart Recommendations", "matching mentors, schemes, incubators"),
    ("02:45", "IPR Filing Tracker", "patent stages & application status"),
    ("03:30", "Govt Integrations", "Aadhaar / DigiLocker / Startup India / ONDC"),
    ("04:00", "Analytics & Impact", "dashboard metrics, sector/district impact"),
    ("04:45", "Audit & Security", "role-based access, audit log, rate limiting"),
]
y = Inches(1.4)
for t, title, desc in steps:
    txt(s, Inches(0.8), y + Inches(0.05), Inches(0.9), Inches(0.4),
        [(t, {'size': 13, 'bold': True, 'color': GREEN})])
    txt(s, Inches(1.8), y + Inches(0.05), Inches(5.0), Inches(0.4),
        [(title, {'size': 14, 'bold': True, 'color': DARK})])
    txt(s, Inches(6.8), y + Inches(0.05), Inches(5.9), Inches(0.4),
        [(desc, {'size': 12, 'color': GRAY})])
    y += Inches(0.58)
footer(s, 8)

# ---------------------------------------------------------------- Slide 9 - Impact
s = slide()
header(s, "Impact & Differentiation", "Why UdaanSetu Wins")
cols = [
    ("Fragmented → Unified", "One platform covering the full research-to-startup journey instead of point solutions."),
    ("Invisible → Discoverable", "Every scheme, mentor, and incubator surfaced to the right innovator automatically."),
    ("Opaque → Transparent", "Clear IPR tracking, risk visibility, and audit trails for all stakeholders."),
    ("Isolated → Connected", "Investors, mentors, incubators, and government all see the same living pipeline."),
]
x = Inches(0.6)
for title, desc in cols:
    rect(s, x, Inches(1.5), Inches(2.9), Inches(3.1), SOFT, line=RGBColor(0xBB, 0xF7, 0xD0))
    txt(s, x + Inches(0.25), Inches(1.8), Inches(2.4), Inches(1.0),
        [(title, {'size': 16, 'bold': True, 'color': MID})])
    txt(s, x + Inches(0.25), Inches(2.7), Inches(2.4), Inches(1.8),
        [(desc, {'size': 12, 'color': GRAY})])
    x += Inches(3.12)
rect(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.3), DARK)
txt(s, Inches(1.0), Inches(5.15), Inches(11.3), Inches(1.0),
    [("Sustainable, scalable, and demo-ready — with a real-data ML training path and production-grade security. "
      "This isn't just a prototype; it's the blueprint for India's innovation highway.", 
      {'size': 15, 'bold': True, 'color': WHITE})], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 9)

# ---------------------------------------------------------------- Slide 10 - Roadmap
s = slide()
header(s, "Roadmap", "Beyond the Hackathon")
rows = [
    ("Phase 1 — Done", "Core platform, 60+ APIs, ML pipeline, govt integrations, tests, Docker"),
    ("Phase 2 — Next", "Ingest real research / patent / scheme datasets; retrain models"),
    ("Phase 3", "Production deployment, CI/CD, monitoring, model drift detection"),
    ("Phase 4", "Govt API live connections, mobile app, multi-language support"),
]
y = Inches(1.5)
for phase, desc in rows:
    rect(s, Inches(0.6), y, Inches(3.4), Inches(0.95), DARK)
    txt(s, Inches(0.75), y + Inches(0.2), Inches(3.1), Inches(0.6),
        [(phase, {'size': 14, 'bold': True, 'color': WHITE})])
    rect(s, Inches(4.1), y, Inches(8.6), Inches(0.95), LIGHT, line=RGBColor(0xE5, 0xE7, 0xEB))
    txt(s, Inches(4.3), y + Inches(0.2), Inches(8.2), Inches(0.6),
        [(desc, {'size': 13, 'color': GRAY})])
    y += Inches(1.1)
footer(s, 10)

# ---------------------------------------------------------------- Slide 11 - Team
s = slide()
header(s, "Team", "Fill Details")
rect(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(4.4), LIGHT, line=RGBColor(0xE5, 0xE7, 0xEB))
txt(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5),
    [("Team Members", {'size': 16, 'bold': True, 'color': DARK})])
members = [
    ("Member 1", "[Name] — Team Leader", "[Email]"),
    ("Member 2", "[Name] — Role", "[Email]"),
    ("Member 3", "[Name] — Role", "[Email]"),
    ("Member 4", "[Name] — Role", "[Email]"),
    ("Member 5", "[Name] — Role", "[Email]"),
    ("Member 6", "[Name] — Role", "[Email]"),
]
y = Inches(2.2)
for i, (ph, name, email) in enumerate(members):
    if i == 3:
        y = Inches(2.2); 
    rect(s, Inches(0.9) + (Inches(3.95) * (i % 3)), y, Inches(3.7), Inches(1.1), WHITE, line=RGBColor(0xE5, 0xE7, 0xEB))
    txt(s, Inches(1.1) + (Inches(3.95) * (i % 3)), y + Inches(0.12), Inches(3.3), Inches(0.4),
        [(name, {'size': 13, 'bold': True, 'color': DARK})])
    txt(s, Inches(1.1) + (Inches(3.95) * (i % 3)), y + Inches(0.5), Inches(3.3), Inches(0.5),
        [(ph + "  •  " + email, {'size': 11, 'color': GRAY})])
    if i == 2:
        y = Inches(3.5)
txt(s, Inches(0.6), Inches(6.0), Inches(12), Inches(0.5),
    [("Institute: [INSTITUTE NAME]  •  City/State: [CITY, STATE]", {'size': 14, 'bold': True, 'color': MID})])
footer(s, 11)

# ---------------------------------------------------------------- Slide 12 - Thanks
s = slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, Inches(4.3), SW, Pt(3), GREEN)
txt(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.2),
    [("Thank You", {'size': 64, 'bold': True, 'color': WHITE, 'align': PP_ALIGN.CENTER})])
txt(s, Inches(0.9), Inches(3.5), Inches(11.5), Inches(0.8),
    [("UdaanSetu — Bridge to Flight  •  SIH 2026  •  Problem ID SIH1608",
      {'size': 20, 'color': RGBColor(0x86, 0xEF, 0xAC), 'align': PP_ALIGN.CENTER})])
txt(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(0.6),
    [("Questions welcome.", {'size': 18, 'color': RGBColor(0xD1, 0xFA, 0xE5), 'align': PP_ALIGN.CENTER})])

prs.save(r"C:\Users\Rudra\Desktop\UdaanSetu\SIH_Submission\UdaanSetu_SIH2026_Demo.pptx")
print("PPT saved: SIH_Submission/UdaanSetu_SIH2026_Demo.pptx")